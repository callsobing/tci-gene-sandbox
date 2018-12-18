from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.shapes.shapetree import SlideShapes
from pptx.util import Inches
import sys
import datetime

now = datetime.datetime.now().strftime("%Y-%m-%d %H:%M")
user_id = sys.argv[1]
uuid = sys.argv[2]
prs = Presentation('data/pptx_template.pptx')

file_fh = open("reports/%s/%s/selected_figures.txt" % (user_id, uuid))

first = True
description = ""
gene_platform_list = {}
for line in file_fh:
    line = line.rstrip()
    if first:
        description = line
        first = False
        continue
file_fh.close()

first_slide = prs.slides[0]
for shape in first_slide.shapes:
    if shape.has_text_frame:
        if shape.text == "title":
            shape.text = description
        elif shape.text == "date":
            shape.text = now

file_fh = open("reports/%s/%s/selected_figures.txt" % (user_id, uuid))

first = True
for line in file_fh:
    line = line.rstrip()
    if first:
        first = False
        continue
    pairs = line.split("###")
    gene = pairs[0]
    platform = pairs[1]
    # 新增slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)

    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = platform

    top = Inches(1.4)
    left = Inches(2.1)
    height = Inches(3.5)
    pic = slide.shapes.add_picture("reports/%s/%s/%s.png" % (user_id, uuid, gene), left, top, height=height)
file_fh.close()

prs.save('reports/%s/%s/%s.pptx' % (user_id, uuid, uuid))
file_fh.close()
