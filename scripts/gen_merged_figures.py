#!/usr/bin/python
import json
import pandas as pd
import matplotlib
matplotlib.use('Agg')
from pylab import *
import numpy as np
from scipy import stats
import os
import sys
plt.rcParams['font.sans-serif'] = ['SimHei']
import datetime
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.shapes.shapetree import SlideShapes
from pptx.util import Inches

now = datetime.datetime.now().strftime("%Y-%m-%d %H:%M")
prs = Presentation('data/pptx_template.pptx')


file_name = sys.argv[1]
samples_fh = sys.argv[2]
user_id = sys.argv[3]
report_uuid = sys.argv[4]

file_fh = open("reports/%s/%s/selected_figures.txt" % (user_id, report_uuid), encoding="utf-8")
platform_genes = {}
selected_genes = []
first = True
for line in file_fh:
    line = line.rstrip()
    if first:
        first = False
        continue
    pairs = line.split("###")
    gene = pairs[0]
    platform = pairs[1]
    selected_genes.append(gene)
    if platform not in platform_genes:
        platform_genes[platform] = []
    platform_genes[platform].append(gene)
file_fh.close()

neg_genes = ["NEG_A", "NEG_B", "NEG_C", "NEG_D", "NEG_E", "NEG_F", "NEG_G", "NEG_H"]

platforms = {
            "抗氧化": {"up": ["SOD1","SOD2","GPX1","CAT"], "down": []},
            "抗老": {"up": ["CCT2","CCT5","CCT6A","CCT7","CCT8","Pink1","Parkin","Atg1","Atg8","SIRT1","FOXO","NADSYN","MRPS5","Ubl-5","SOD3"], "down": ["PARP1","PARP2"]},
            "DNA修復": {"up": ["UNG","OGG1","MPG","APEX1","ERCC1","ERCC6","XPA","XRCC1","XRCC5","MSH2","MLH1","MSH6"], "down": []},
            "免疫": {"up": ["IL-1B","IL-8","IL-6","IL-10","IL-18","TNF-a"], "down": []},
            "護胃": {"up": ["SOD1","SOD2","GPX1"], "down": ["IL-1B","IL-6","IL-8"]},
            "護眼": {"up": [], "down": ["VEGFA","CASP3","CASP8","IL-1B","IL-8","TNF-a"]},
            "抗黑色素生成": {"up": [], "down": ["TYR","TYRP1","MC1R","MITF"]},
            "膠原蛋白合成-組合-降解": {"up": ["COL1A1","COL1A2","COL4A1","COL4A4","COL4A5","TIMP1","ELN","FBN1","LOX","HAS2","HAS3"], "down": ["MMP1","MMP9","MMP2"]},
            "抗發炎": {"up": ["IL-10","TGFB","IL4"], "down": ["IL-1B","IL-8","IL-6","IL-18","TNF-a","IL-16","IL23","IL12A","IFNG","IL3"]},
            "細胞凋亡": {"up": ["BCL-2"], "down": ["BAX","BCLXL","BAD","CASP9","AIFM1","EndoG"]},
            "心血管保健": {"up": ["PTGIS","NOS3","PLAT","PROC"], "down": ["EDN1","VWF","F3","SERPINE1","PDGFC","FGF2","IGF2BP3","IGF1R","IL-8","IL-6","ICAM1","VCAM1","CASP8"]},
            "晝夜節律": {"up": ["SIRT1","CLOCK","BMAL1 (ARNTL)","PER2","CRY","KPNB1"], "down": []},
            "LPS模擬體內發炎": {"up": ["NOS3"], "down": ["ICAM1","VCAM1","IL-8"]},
            "非酒精性肝損傷": {"up": ["UNG","OGG1","MPG","APEX1","ERCC1","ERCC6","XPA","XRCC1","XRCC5","MSH2","MLH1","MSH6","SOD1","SOD2","GPX1","CAT"], "down": []},
            "皮膚角質保濕": {"up": ["Tgm1","Krt1","Keratin 10","Keratin 14","AQP3","FLG-F","SMPD1","GBA","HAS2","HAS3"], "down": []},
            "脂肪肝": {"up": ["PPAR-g","PPAR-a"], "down": ["SREBP-1c (SREBF1)","SCD1 (SCD)","ACC (ACACA)"]},
            "提升HDL": {"up": ["CETP","SCARB1","apoA-I (APA1)","LDLR","ABCA1"], "down": []},
            "健髮平台": {"up": ["KROX20","SCF","VEGFA","IGF1"], "down": ["SRD5A1","SRD5A2","AR","TGFB","BDNF"]},
            "端粒酶活性平台": {"up": ["TERT","TERC","RTEL1"], "down": []},
            "免疫活化與分化": {"up": [], "down": ["CD40","ERBB2","LIF","MALT1","NCK1","PAF1","DYNLL2","GRK5","PSMD4","RDH10","RELB","SCARF1","TNFSF14","ABR","IL13","IL4R","IL5RA","RELA"]}
            }

platform_map = {
            "抗氧化": "platform1",
            "抗老": "platform2",
            "DNA修復":  "platform3",
            "免疫":  "platform4",
            "護胃":  "platform5",
            "護眼":  "platform6",
            "抗黑色素生成":  "platform7",
            "膠原蛋白合成-組合-降解":  "platform8",
            "抗發炎":  "platform9",
            "細胞凋亡":  "platform10",
            "心血管保健":  "platform11",
            "晝夜節律":  "platform12",
            "LPS模擬體內發炎":  "platform13",
            "非酒精性肝損傷":  "platform14",
            "皮膚角質保濕":  "platform15",
            "脂肪肝":  "platform16",
            "提升HDL":  "platform17",
            "健髮平台":  "platform18",
            "端粒酶活性平台":  "platform19",
            "免疫活化與分化":  "platform20"
}

def significance_score(mean1, mean2):
    try:
        (statistic, pvalue) = stats.ttest_ind(mean1, mean2)
    except:
        pvalue = 1
    if pvalue < 0.001:
        return 3.
    elif pvalue < 0.01:
        return 2.
    elif pvalue < 0.05:
        return 1.
    return 0


def expression_direction_score(direction, expression):
    difference = average(expression) - 1
    check_direction = difference * direction
    multiplier = 1
    if check_direction < 0: # 和預期方向不同就扣分
        multiplier = -1

    if abs(difference) > 0.5:
        return 3. * multiplier
    elif abs(difference) > 0.3:
        return 2. * multiplier
    elif abs(difference) > 0.1:
        return 1. * multiplier
    return 0


def create_directory(directory_path):
    if not os.path.exists(directory_path):
        os.makedirs(directory_path)
        return True
    else:
        return False


def label_significance(pos_in_fig, mean1, mean2, std, ymax, N):
    try:
        (statistic, pvalue) = stats.ttest_ind(mean1, mean2)
    except:
        pvalue = 1
    text = ""
    if pvalue < 0.001:
        text = "***"
    elif pvalue < 0.01:
        text = "**"
    elif pvalue < 0.05:
        text = "*"
    if text:
        # Annotate significance level
        plt.annotate(text, xy=(pos_in_fig, average(mean2) + std + ymax * 0.005), fontsize=(15 - N*0.7))
        # Annotate Relative Expression ratio
    # plt.annotate("{:.2f}".format(average(mean2)), xy=(pos_in_fig - 0.05, average(mean2) - ymax * 0.05), fontsize=12)


def plot_platform(gene_details_map, platform_name, user_id, report_uuid, sample_ids):
    mock_means = []
    c1t1_means = []
    c1t2_means = []
    c2t1_means = []
    c2t2_means = []
    gene_names = []
    mock_errors = [[], []]
    c1t1_errors = [[], []]
    c1t2_errors = [[], []]
    c2t1_errors = [[], []]
    c2t2_errors = [[], []]
    N = len(platform_genes[platform_name])  # Number of genes
    xpos = np.arange(N)
    ind = []
    for i in xpos:
        ind.append(xpos[i] + 0.2)
    ind = np.array(ind)
    width = 0.18  # the width of the bars
    fig, ax = plt.subplots()

    for gene in platform_genes[platform_name]:
        gene_names.append(gene)
        mock_mean = average(gene_details_map[gene]["mock"]["fold_change"])
        mean_c1t1 = average(gene_details_map[gene]["c1t1"]["fold_change"])
        mean_c1t2 = average(gene_details_map[gene]["c1t2"]["fold_change"])
        mean_c2t1 = average(gene_details_map[gene]["c2t1"]["fold_change"])
        mean_c2t2 = average(gene_details_map[gene]["c2t2"]["fold_change"])

        mock_means.append(mock_mean)
        c1t1_means.append(mean_c1t1)
        c1t2_means.append(mean_c1t2)
        c2t1_means.append(mean_c2t1)
        c2t2_means.append(mean_c2t2)

        mock_errors[0].append(0)
        mock_errors[1].append(gene_details_map[gene]["mock"]["std"]/2)
        c1t1_errors[0].append(0)
        c1t1_errors[1].append(gene_details_map[gene]["c1t1"]["std"] / 2)
        c1t2_errors[0].append(0)
        c1t2_errors[1].append(gene_details_map[gene]["c1t2"]["std"] / 2)
        c2t1_errors[0].append(0)
        c2t1_errors[1].append(gene_details_map[gene]["c2t1"]["std"] / 2)
        c2t2_errors[0].append(0)
        c2t2_errors[1].append(gene_details_map[gene]["c2t2"]["std"] / 2)

    labels = ['控制組', sample_ids[0], sample_ids[1], sample_ids[2], sample_ids[3]]

    plt.bar(ind, mock_means, width, color="#3D3A4B", align='center', linewidth=0, label=labels[0])
    plt.bar(ind + width, c1t1_means, width, color="#705D56", align='center', linewidth=0, label=labels[1])
    plt.bar(ind + width * 2, c1t2_means, width, color="#B19994", align='center', linewidth=0, label=labels[2])
    plt.bar(ind + width * 3, c2t1_means, width, color="#D3C0CD", align='center', linewidth=0, label=labels[3])
    plt.bar(ind + width * 4, c2t2_means, width, color="#E3DFFF", align='center', linewidth=0, label=labels[4])

    plotline2, caplines2, barlinecols2 = ax.errorbar(ind + width, c1t1_means, yerr=c1t1_errors, lolims=True, ls='None', color='black', barsabove=True)
    plotline3, caplines3, barlinecols3 = ax.errorbar(ind + width * 2, c1t2_means, yerr=c1t2_errors, lolims=True, ls='None', color='black', barsabove=True)
    plotline4, caplines4, barlinecols4 = ax.errorbar(ind + width * 3, c2t1_means, yerr=c2t1_errors, lolims=True, ls='None', color='black', barsabove=True)
    plotline5, caplines5, barlinecols5 = ax.errorbar(ind + width * 4, c2t2_means, yerr=c2t2_errors, lolims=True, ls='None', color='black', barsabove=True)
    lgd = ax.legend(bbox_to_anchor=(1.05, 1), loc='upper left', borderaxespad=0.)

    caplines2[0].set_marker('.')
    caplines2[2].set_marker('_')
    caplines2[0].set_markersize(0)
    caplines2[2].set_markersize(15 - N)

    caplines3[0].set_marker('.')
    caplines3[2].set_marker('_')
    caplines3[0].set_markersize(0)
    caplines3[2].set_markersize(15 - N)

    caplines4[0].set_marker('.')
    caplines4[2].set_marker('_')
    caplines4[0].set_markersize(0)
    caplines4[2].set_markersize(15 - N)

    caplines5[0].set_marker('.')
    caplines5[2].set_marker('_')
    caplines5[0].set_markersize(0)
    caplines5[2].set_markersize(15 - N)

    ymax = max([max(mock_means), max(c1t1_means), max(c1t2_means), max(c2t1_means), max(c2t2_means)]) + max([max(c1t1_errors[1]), max(c1t2_errors[1]), max(c2t1_errors[1]), max(c2t2_errors[1])]) * 2
    plt.ylim(0, ymax)
    plt.ylabel('Relative Expression Ratio', fontsize="x-large")
    ax.yaxis.set_ticks_position('left')
    ax.xaxis.set_ticks_position('bottom')
    ax.set_xticks(ind + width * 2)
    ax.set_xticklabels(platform_genes[platform_name])
    # plt.xticks(x_pos, labels, color='k', fontsize="x-large")
    plt.yticks(fontsize="x-large")
    plt.gca().spines['right'].set_visible(False)
    plt.gca().spines['top'].set_visible(False)

    for i in range(len(platform_genes[platform_name])):
        gene = platform_genes[platform_name][i]
        label_significance(ind[i]+width*(0+1*0.8), gene_details_map[gene]["mock"]["fold_change"], gene_details_map[gene]["c1t1"]["fold_change"], c1t1_errors[1][i], ymax, N)
        label_significance(ind[i]+width*(1+1*0.8), gene_details_map[gene]["mock"]["fold_change"], gene_details_map[gene]["c1t2"]["fold_change"], c1t2_errors[1][i], ymax, N)
        label_significance(ind[i]+width*(2+1*0.8), gene_details_map[gene]["mock"]["fold_change"], gene_details_map[gene]["c2t1"]["fold_change"], c2t1_errors[1][i], ymax, N)
        label_significance(ind[i]+width*(3+1*0.8), gene_details_map[gene]["mock"]["fold_change"], gene_details_map[gene]["c2t2"]["fold_change"], c2t2_errors[1][i], ymax, N)
    file_name_path = os.path.normcase("reports/%s/%s/%s.png" % (user_id, report_uuid, platform_map[platform_name]))
    fig.savefig(file_name_path, bbox_extra_artists=(lgd,), bbox_inches='tight')
    plt.cla()
    plt.close(fig)


selected_sample_fh = open(samples_fh, encoding="utf-8")
mock_samples = []
c1t1_samples = []
c1t2_samples = []
c2t1_samples = []
c2t2_samples = []
sample_identifiers = []
count = 0
for line in selected_sample_fh:
    line = line.rstrip()
    splitted = line.split("\t")
    if count == 0:
        sample_identifiers = splitted
        count += 1
        continue
    if count == 1:
        mock_samples = splitted
        count += 1
        continue
    if count == 2:
        c1t1_samples = splitted
        count += 1
        continue
    if count == 3:
        c1t2_samples = splitted
        count += 1
        continue
    if count == 4:
        c2t1_samples = splitted
        count += 1
        continue
    if count == 5:
        c2t2_samples = splitted
        continue
selected_sample_fh.close()

with open(file_name, 'r', errors="ignore", encoding="utf-8") as f:
    nCounter_fh = f.readlines()

samples_idx = {}
gene_names = []
expression_map = {}
count = 0
for line in nCounter_fh:
    line = line.rstrip()
    splitted = line.split("\t")
    if count == 1:
        for idx in range(8, len(splitted)):
            samples_idx["{:02d}".format(idx) + ":" + splitted[idx]] = idx
            expression_map["{:02d}".format(idx)  + ":" + splitted[idx]] = []
        count += 1
        continue
    if count < 3:
        count += 1
        continue
    gene_names.append(splitted[0])
    for sample in samples_idx:
        expression_map[sample].append(splitted[samples_idx[sample]])  # 最後輸出格式: expression_map: 以sample名字為key，按照次序記錄各基因表現量

gene_details_map = {}
# 初始化gene maps
for gene_idx in range(len(gene_names)):
    gene_details_map[gene_names[gene_idx]] = \
        {
        "mock": {"fold_change": [1, 1, 1], "std": 0},
        "c1t1": {"fold_change": [], "std": 0.0},
        "c1t2": {"fold_change": [], "std": 0.0},
        "c2t1": {"fold_change": [], "std": 0.0},
        "c2t2": {"fold_change": [], "std": 0.0}
        }

threshold = 20.0001

for gene_idx in range(len(gene_names)):
    mock_sum = 0
    count = 0
    for mock_id in mock_samples:
        mock_idx = samples_idx[mock_id]
        if float(expression_map[mock_id][gene_idx]) < threshold:
            count += 1
            continue
        mock_sum += float(expression_map[mock_id][gene_idx])
    if count == 2 or mock_sum == 0: # 如果剛好兩個mock沒有超過threshold就丟棄這個gene
        gene_details_map.pop(gene_names[gene_idx])
        continue
    count = 0
    mock_avg = mock_sum / float(len(mock_samples))

    for cond1_sample in c1t1_samples:
        sample_idx = samples_idx[cond1_sample]
        if float(expression_map[cond1_sample][gene_idx]) < threshold:
            count += 1
            continue
        gene_details_map[gene_names[gene_idx]]["c1t1"]["fold_change"].append(float(expression_map[cond1_sample][gene_idx]) / mock_avg)
    gene_details_map[gene_names[gene_idx]]["c1t1"]["std"] = np.std(gene_details_map[gene_names[gene_idx]]["c1t1"]["fold_change"])
    if count > 1: # 如果超過兩個mock沒有超過threshold就丟棄這個gene
        gene_details_map[gene_names[gene_idx]]["c1t1"]["fold_change"] = []
        gene_details_map[gene_names[gene_idx]]["c1t1"]["std"] = 0.0
    count = 0

    for cond2_sample in c1t2_samples:
        sample_idx = samples_idx[cond2_sample]
        if float(expression_map[cond2_sample][gene_idx]) < threshold:
            count += 1
            continue
        gene_details_map[gene_names[gene_idx]]["c1t2"]["fold_change"].append(float(expression_map[cond2_sample][gene_idx]) / mock_avg)
    gene_details_map[gene_names[gene_idx]]["c1t2"]["std"] = np.std(gene_details_map[gene_names[gene_idx]]["c1t2"]["fold_change"])
    if count > 1:  # 如果超過兩個mock沒有超過threshold就丟棄這個gene
        gene_details_map[gene_names[gene_idx]]["c1t2"]["fold_change"] = []
        gene_details_map[gene_names[gene_idx]]["c1t2"]["std"] = 0.0
    count = 0

    for cond1_sample in c2t1_samples:
        sample_idx = samples_idx[cond1_sample]
        if float(expression_map[cond1_sample][gene_idx]) < threshold:
            count += 1
            continue
        gene_details_map[gene_names[gene_idx]]["c2t1"]["fold_change"].append(float(expression_map[cond1_sample][gene_idx]) / mock_avg)
    gene_details_map[gene_names[gene_idx]]["c2t1"]["std"] = np.std(gene_details_map[gene_names[gene_idx]]["c2t1"]["fold_change"])
    if count > 1:  # 如果超過兩個mock沒有超過threshold就丟棄這個gene
        gene_details_map[gene_names[gene_idx]]["c2t1"]["fold_change"] = []
        gene_details_map[gene_names[gene_idx]]["c2t1"]["std"] = 0.0
    count = 0

    for cond2_sample in c2t2_samples:
        sample_idx = samples_idx[cond2_sample]
        if float(expression_map[cond2_sample][gene_idx]) < threshold:
            count += 1
            continue
        gene_details_map[gene_names[gene_idx]]["c2t2"]["fold_change"].append(float(expression_map[cond2_sample][gene_idx]) / mock_avg)
    gene_details_map[gene_names[gene_idx]]["c2t2"]["std"] = np.std(gene_details_map[gene_names[gene_idx]]["c2t2"]["fold_change"])
    if count > 1:  # 如果超過兩個mock沒有超過threshold就丟棄這個gene
        gene_details_map[gene_names[gene_idx]]["c2t2"]["fold_change"] = []
        gene_details_map[gene_names[gene_idx]]["c2t2"]["std"] = 0.0
    count = 0

file_fh = open("reports/%s/%s/selected_figures.txt" % (user_id, report_uuid), encoding="utf-8")

first = True
description = ""
gene_platform_list = {}
for line in file_fh:
    line = line.rstrip()
    if first:
        description = line
        first = False
        continue
file_fh.close()

first_slide = prs.slides[0]
for shape in first_slide.shapes:
    if shape.has_text_frame:
        if shape.text == "title":
            shape.text = description
        elif shape.text == "date":
            shape.text = now

# 這邊應該要新增一個資料夾叫做report專門存相關資料
# report下面新開uuid的資料夾 - "/使用者名稱/uuid/"
for platform_name in platform_genes:
    plot_platform(gene_details_map, platform_name, user_id, report_uuid, sample_identifiers)
    # 新增slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)

    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = platform_name

    top = Inches(1.5)
    left = Inches(2.4)
    height = Inches(3.5)
    pic = slide.shapes.add_picture("reports/%s/%s/%s.png" % (user_id, report_uuid, platform_map[platform_name]), left, top, height=height)

prs.save('reports/%s/%s/%s.pptx' % (user_id, report_uuid, report_uuid))